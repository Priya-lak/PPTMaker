import os

from pptx.enum.text import PP_ALIGN
from pptx.util import Inches


class ImageHandler:
    """Handles image-related slide operations"""

    def __init__(self, style):
        self.style = style

    def add_image_slide(
        self,
        slide_manager,
        title,
        image_path,
        position=None,
        caption="",
        layout_type="standard",
    ):
        """Create a slide with a title, image, and optional caption"""
        slide = slide_manager.create_blank_slide_with_title(title)

        if not os.path.exists(image_path):
            self._add_image_not_found_placeholder(slide, image_path)
            return slide

        return self._apply_image_layout(
            slide, image_path, caption, layout_type, position
        )

    def _add_image_not_found_placeholder(self, slide, image_path):
        """Add a placeholder when image is not found"""
        placeholder = slide.shapes.add_textbox(
            Inches(2), Inches(3), Inches(9.33), Inches(2)
        )
        placeholder.text = f"Image not found: {image_path}"
        placeholder_paragraph = placeholder.text_frame.paragraphs[0]
        self.style.apply_font_style(placeholder_paragraph, "body_medium")

    def _apply_image_layout(
        self, slide, image_path, caption, layout_type, position=None
    ):
        """Apply specific image layout based on layout_type"""
        if layout_type == "side_by_side":
            return self._create_side_by_side_layout(slide, image_path, caption)
        elif layout_type == "image_bottom":
            return self._create_image_bottom_layout(slide, image_path, caption)
        elif layout_type == "image_title_caption":
            return self._create_image_title_caption_layout(slide, image_path, caption)
        else:
            return self._create_standard_layout(slide, image_path, caption, position)

    def _create_side_by_side_layout(self, slide, image_path, caption):
        """Image left, caption right"""
        slide.shapes.add_picture(
            image_path, Inches(0.5), Inches(1.5), width=Inches(5.5), height=Inches(4.5)
        )
        if caption:
            caption_box = slide.shapes.add_textbox(
                Inches(6.2), Inches(1.5), Inches(5.5), Inches(4.5)
            )
            self._format_caption(caption_box, caption)
        return slide

    def _create_image_bottom_layout(self, slide, image_path, caption):
        """Caption above, image below"""
        if caption:
            caption_box = slide.shapes.add_textbox(
                Inches(1), Inches(1.5), Inches(10), Inches(2)
            )
            self._format_caption(caption_box, caption)
        slide.shapes.add_picture(
            image_path, Inches(1), Inches(3.5), width=Inches(10), height=Inches(3.5)
        )
        return slide

    def _create_image_title_caption_layout(self, slide, image_path, caption):
        """Image > Title > Caption (bullets and description)"""
        slide.shapes.add_picture(
            image_path, Inches(1), Inches(1.2), width=Inches(10), height=Inches(3)
        )
        if caption:
            caption_box = slide.shapes.add_textbox(
                Inches(1), Inches(4.3), Inches(10), Inches(2.5)
            )
            caption_frame = caption_box.text_frame
            caption_frame.wrap = True
            for line in caption.split("\n"):
                p = caption_frame.add_paragraph()
                p.text = line
                p.level = 0
                self.style.apply_font_style(p, "body_medium")
        return slide

    def _create_standard_layout(self, slide, image_path, caption, position):
        """Default standard layout (image centered, caption below)"""
        slide.shapes.add_picture(
            image_path,
            Inches(position.left),
            Inches(position.top),
            width=Inches(position.width),
            height=Inches(position.height),
        )
        if caption:
            caption_box = slide.shapes.add_textbox(
                Inches(2), Inches(6.5), Inches(9.33), Inches(0.5)
            )
            self._format_caption(caption_box, caption)
        return slide

    def _format_caption(self, caption_box, caption):
        """Format caption text consistently"""
        caption_frame = caption_box.text_frame
        caption_frame.clear()
        caption_frame.word_wrap = True
        caption_frame.text = caption
        caption_paragraph = caption_frame.paragraphs[0]
        self.style.apply_font_style(caption_paragraph, "caption")
        caption_paragraph.alignment = PP_ALIGN.LEFT
