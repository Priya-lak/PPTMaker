from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches


class ChartHandler:
    """Handles chart creation and formatting"""

    def __init__(self, style):
        self.style = style

    def add_chart_slide(self, slide_manager, title, chart_data, chart_type="column"):
        """Create a slide with a chart"""
        slide = slide_manager.create_blank_slide_with_title(title)

        # Update title formatting for charts
        title_shapes = [shape for shape in slide.shapes if hasattr(shape, "text_frame")]
        if title_shapes:
            title_paragraph = title_shapes[0].text_frame.paragraphs[0]
            self.style.apply_font_style(title_paragraph, "title_small")

        # Create and add chart
        chart_data_obj = self._prepare_chart_data(chart_data)
        chart_type_enum = self._get_chart_type_enum(chart_type)

        x, y, cx, cy = Inches(1.5), Inches(2), Inches(10), Inches(4.5)
        graphic_frame = slide.shapes.add_chart(
            chart_type_enum, x, y, cx, cy, chart_data_obj
        )

        # Customize chart appearance
        chart = graphic_frame.chart
        chart.has_legend = True
        chart.legend.position = 2  # Right side

        return slide

    def _prepare_chart_data(self, chart_data):
        """Prepare chart data object"""
        chart_data_obj = CategoryChartData()
        categories = list(chart_data.keys())
        values = list(chart_data.values())

        chart_data_obj.categories = categories
        chart_data_obj.add_series("Series 1", values)

        return chart_data_obj

    def _get_chart_type_enum(self, chart_type):
        """Get chart type enumeration"""
        chart_types = {
            "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "pie": XL_CHART_TYPE.PIE,
            "line": XL_CHART_TYPE.LINE,
        }
        return chart_types.get(chart_type.lower(), XL_CHART_TYPE.COLUMN_CLUSTERED)
