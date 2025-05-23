from pptx.util import Inches


class TableHandler:
    """Handles table creation and formatting"""

    def __init__(self, style):
        self.style = style

    def create_table_slide(self, slide_manager, title, table_data):
        """Create a slide with a table"""
        slide = slide_manager.create_blank_slide_with_title(title)

        # Update title formatting for tables
        title_shapes = [shape for shape in slide.shapes if hasattr(shape, "text_frame")]
        if title_shapes:
            title_paragraph = title_shapes[0].text_frame.paragraphs[0]
            self.style.apply_font_style(title_paragraph, "title_small")

        # Create table if data is provided
        if table_data and len(table_data) > 0:
            self._create_and_format_table(slide, table_data)

        return slide

    def _create_and_format_table(self, slide, table_data):
        """Create and format table with data"""
        rows = len(table_data)
        cols = len(table_data[0]) if table_data else 0

        if rows > 0 and cols > 0:
            # Table dimensions
            left = self.style.get_dimension("margin_large")
            top = Inches(2)
            width = Inches(11.33)
            height = Inches(4.5)

            table = slide.shapes.add_table(rows, cols, left, top, width, height).table

            # Populate and format table
            for row_idx, row_data in enumerate(table_data):
                for col_idx, cell_data in enumerate(row_data):
                    cell = table.cell(row_idx, col_idx)
                    cell.text = str(cell_data)

                    cell_paragraph = cell.text_frame.paragraphs[0]

                    if row_idx == 0:  # Header row
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = self.style.shapes[
                            "table_header_fill"
                        ]
                        self.style.apply_font_style(cell_paragraph, "table_header")
                    else:  # Data rows
                        self.style.apply_font_style(cell_paragraph, "body_medium")
