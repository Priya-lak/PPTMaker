from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches


class ShapeHandler:
    """Handles shape and text box creation"""

    def __init__(self, style):
        self.style = style

    def add_text_box(
        self, slide, text, left, top, width, height, font_style="body_large"
    ):
        """Add a styled text box to a slide"""
        text_box = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        text_frame = text_box.text_frame
        text_frame.text = text

        paragraph = text_frame.paragraphs[0]
        self.style.apply_font_style(paragraph, font_style)

        return text_box

    def add_shape(self, slide, shape_type, left, top, width, height, fill_color=None):
        """Add a shape to a slide"""
        shape_enum = self._get_shape_enum(shape_type)

        shape = slide.shapes.add_shape(
            shape_enum, Inches(left), Inches(top), Inches(width), Inches(height)
        )

        self._apply_shape_fill(shape, fill_color)
        return shape

    def _get_shape_enum(self, shape_type):
        """Get shape enumeration based on type"""
        shape_types = {
            "rectangle": MSO_SHAPE.RECTANGLE,
            "circle": MSO_SHAPE.OVAL,
            "oval": MSO_SHAPE.OVAL,
        }
        return shape_types.get(shape_type.lower(), MSO_SHAPE.RECTANGLE)

    def _apply_shape_fill(self, shape, fill_color):
        """Apply fill color to shape"""
        fill = shape.fill
        fill.solid()

        if fill_color:
            if isinstance(fill_color, tuple) and len(fill_color) == 3:
                fill.fore_color.rgb = RGBColor(*fill_color)
            elif isinstance(fill_color, str):
                fill.fore_color.rgb = self.style.get_color(fill_color)
        else:
            fill.fore_color.rgb = self.style.get_color("primary_light")
