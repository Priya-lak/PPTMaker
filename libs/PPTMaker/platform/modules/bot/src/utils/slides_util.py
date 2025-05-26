from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.slide import Slide
from pptx.util import Inches

from libs.PPTMaker.platform.modules.bot.src.models.slide_models import SlideLayout
from libs.PPTMaker.platform.modules.bot.src.utils.styles.base_style import (
    BasePresentationStyle,
)


class SlideBuilder:
    """Unified slide creation and management"""

    def __init__(self, presentation: Presentation, style: BasePresentationStyle):
        self.prs = presentation
        self.style = style

    def create_slide(self, layout: str, **kwargs) -> Slide:
        """Create a slide with the specified layout and content"""
        method_map = {
            "TITLE": self.title_slide,
            "TITLE_AND_CONTENT": self.title_and_content_slide,
            "SECTION_HEADER": self.section_header_slide,
            "TWO_CONTENT": self.two_content_slide,
            "COMPARISON": self.comparison_slide,
            "TITLE_ONLY": self.title_only_slide,
            "BLANK": self.blank_slide_with_title,
            "CONTENT_WITH_CAPTION": self.content_with_caption_slide,
            "PICTURE_WITH_CAPTION": self.picture_with_caption_slide,
        }

        if layout not in method_map:
            raise ValueError(f"Unsupported layout: {layout}")

        return method_map[layout](**kwargs)

    # Core slide creation methods
    def title_slide(self, title: str, subtitle: str = "", **kwargs) -> Slide:
        slide = self._add_slide(SlideLayout.TITLE)
        self._set_title(slide.shapes.title, title, "title_large")
        if subtitle:
            self._set_subtitle(slide.placeholders[1], subtitle)
        return slide

    def title_and_content_slide(self, title: str, points: List[str], **kwargs) -> Slide:
        slide = self._add_slide(SlideLayout.TITLE_AND_CONTENT)
        self._set_title(slide.shapes.title, title, "title_medium")
        self._add_bullet_points(slide.placeholders[1], points)
        return slide

    def section_header_slide(self, title: str, subtitle: str = "", **kwargs) -> Slide:
        slide = self._add_slide(SlideLayout.SECTION_HEADER)
        self._set_title(slide.shapes.title, title, "title_medium")
        if subtitle:
            self._set_subtitle(slide.placeholders[1], subtitle)
        return slide

    def two_content_slide(
        self, title: str, left_points: List[str], right_points: List[str], **kwargs
    ) -> Slide:
        slide = self._add_slide(SlideLayout.TWO_CONTENT)
        self._set_title(slide.shapes.title, title, "title_medium")
        self._add_bullet_points(slide.placeholders[1], left_points)
        self._add_bullet_points(slide.placeholders[2], right_points)
        return slide

    def comparison_slide(
        self,
        title: str,
        left_heading: str,
        left_points: List[str],
        right_heading: str,
        right_points: List[str],
    ) -> Slide:
        slide = self._add_slide(SlideLayout.COMPARISON)
        self._set_title(slide.shapes.title, title, "title_medium")

        # Set headings
        slide.placeholders[1].text = left_heading
        slide.placeholders[3].text = right_heading

        # Add content
        self._add_bullet_points(slide.placeholders[2], left_points)
        self._add_bullet_points(slide.placeholders[4], right_points)
        return slide

    def title_only_slide(self, title: str, **kwargs) -> Slide:
        slide = self._add_slide(SlideLayout.TITLE_ONLY)
        self._set_title(slide.shapes.title, title, "title_medium")
        return slide

    def blank_slide_with_title(self, title: str, **kwargs) -> Slide:
        slide = self._add_slide(SlideLayout.BLANK)
        title_box = slide.shapes.add_textbox(
            self.style.get_dimension("margin_standard"),
            Inches(0.3),
            Inches(12),
            self.style.get_dimension("spacing_large"),
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        # self.style.apply_font_style(title_frame.paragraphs[0], "title_medium")
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        return slide

    def content_with_caption_slide(
        self, title: str, points: List[str], **kwargs
    ) -> Slide:
        slide = self._add_slide(SlideLayout.CONTENT_WITH_CAPTION)
        slide.placeholders[0].text_frame.clear()
        self._set_title(slide.shapes.title, title, "title_medium")
        self._add_bullet_points(slide.placeholders[1], points)
        return slide

    def picture_with_caption_slide(
        self, title: str, image_path: str, caption: str, **kwargs
    ) -> Slide:
        slide = self._add_slide(SlideLayout.PICTURE_WITH_CAPTION)
        self._set_title(slide.shapes.title, title, "title_medium")

        # Insert image and caption
        slide.placeholders[1].insert_picture(image_path)
        slide.placeholders[2].text = caption
        # self.style.apply_font_style(slide.placeholders[2].text_frame.paragraphs[0], "caption")
        return slide

    # Helper methods
    def _add_slide(self, layout_type: SlideLayout) -> Slide:
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[layout_type])
        # self.style.apply_background(slide)
        return slide

    def _set_title(
        self, shape: Any, text: str, style_key: str, align=PP_ALIGN.LEFT
    ) -> None:
        shape.text = text
        p = shape.text_frame.paragraphs[0]
        # self.style.apply_font_style(p, style_key)
        p.alignment = align

    def _set_subtitle(self, placeholder: Any, text: str) -> None:
        placeholder.text = text
        # self.style.apply_font_style(placeholder.text_frame.paragraphs[0], "subtitle")

    def _add_bullet_points(self, shape: Any, points: List[str]) -> None:
        if not points:
            return

        text_frame = shape.text_frame
        text_frame.wrap = True
        text_frame.clear()

        for i, point in enumerate(points):
            p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
            p.text = point
            p.level = 0
            # self.style.apply_font_style(p, "body_large")


# Usage examples:
# builder = SlideBuilder(presentation, style)

# Simple method calls (cleaner for known layouts)
# slide1 = builder.title_slide("My Title", "My Subtitle")
# slide2 = builder.title_and_content_slide("Content", ["Point 1", "Point 2"])

# Dynamic dispatch (useful when layout comes from data)
# slide3 = builder.create_slide(SlideLayout.COMPARISON,
#                              title="Compare",
#                              left_heading="Pro", left_points=["Good"],
#                              right_heading="Con", right_points=["Bad"])
