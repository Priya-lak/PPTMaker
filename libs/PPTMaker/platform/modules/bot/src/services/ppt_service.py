from pprint import pformat
from typing import Dict

from pptx import Presentation
from pptx.slide import Slide

from libs.PPTMaker.platform.modules.bot.src.models.slide_layout_models_2 import (
    LayoutDefinition,
    PresentationResponse,
)
from libs.utils.common.custom_logger import CustomLogger

log = CustomLogger("PPTMaker", is_request=False)

logger, listener = log.get_logger()

listener.start()


class PPT:
    def __init__(self, theme):
        self.prs = Presentation(theme)

    def get_theme_layouts(self):
        layouts = []
        for slide_layout in self.prs.slide_layouts:
            layout_blueprint = LayoutDefinition(
                layout=slide_layout.name,
                placeholders=[
                    {"idx": idx, "name": placeholder.name}
                    for idx, placeholder in enumerate(slide_layout.placeholders)
                ],
            )
            layouts.append(layout_blueprint.model_dump())

        return layouts

    def add_slide(self, layout_identifier):
        """Add slide by layout name or index"""
        if isinstance(layout_identifier, str):
            # Find layout by name
            for i, layout in enumerate(self.prs.slide_layouts):
                if layout.name == layout_identifier:
                    # logger.info(f"adding slide {layout.name}")
                    slide = self.prs.slides.add_slide(self.prs.slide_layouts[i])
                    return slide
            raise ValueError(f"Layout '{layout_identifier}' not found")
        else:
            # Use as index
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[layout_identifier])
            return slide

    def populate_presentation(self, presentation_content: PresentationResponse):
        """Populate presentation with structured content"""
        logger.info("populating ppt")
        for slide_content in presentation_content.slides:
            try:
                # Add slide with the specified layout
                slide = self.add_slide(slide_content.layout)

                # Populate placeholders with content
                self.populate_slide_content(
                    slide, slide_content.placeholders, slide_content.layout
                )

            except ValueError as e:
                logger.info(
                    f"Error adding slide with layout '{slide_content.layout}': {e}"
                )
                continue
            except Exception as e:
                logger.info(f"Unexpected error processing slide: {e}")
                raise

    def populate_slide_content(
        self, slide: Slide, placeholder_content: Dict[str, str], sl_layout: str
    ):
        """Fill slide placeholders with provided content"""

        # Method 1: Direct iteration over slide placeholders
        original_slide_placeholders = [
            {"idx": idx, "name": pl.name} for idx, pl in enumerate(slide.placeholders)
        ]
        logger.debug(
            f"layout: {slide.slide_layout.name}, original slide placeholders {pformat(original_slide_placeholders)}"
        )
        logger.debug(f"layout: {sl_layout}, generated {pformat(placeholder_content)}")
        for idx, placeholder in enumerate(slide.placeholders):
            try:
                placeholder_name = placeholder.name
                content = placeholder_content.get(str(idx))
                if content is None:
                    continue

                # Check if placeholder can accept text
                if placeholder.has_text_frame:
                    if placeholder.text_frame is not None:
                        # Clear existing content
                        placeholder.text_frame.clear()

                        # Split content by newlines
                        lines = content.split("\\n")

                        # Add first line to the default paragraph
                        if lines:
                            placeholder.text_frame.text = lines[0]

                        # Add remaining lines as new paragraphs
                        for line in lines[1:]:
                            p = placeholder.text_frame.add_paragraph()
                            p.text = line

                        # logger.info(
                        #     f"Populated placeholder '{placeholder_name}' with content: {content}"
                        # )
                    else:
                        logger.warning(
                            f"Placeholder '{placeholder_name}' has text frame but frame is None"
                        )
                else:
                    logger.info(
                        f"Placeholder '{placeholder_name}' cannot accept text (no text frame)"
                    )

            except Exception as e:
                logger.error(
                    f"Error processing placeholder '{getattr(placeholder, 'name', 'unknown')}': {e}"
                )
                continue

    def get_available_layouts_json(self):
        """Get available layouts in JSON format for AI prompt"""
        layouts = []
        for slide_layout in self.prs.slide_layouts:
            layouts.append(
                {
                    "name": slide_layout.name,
                    "placeholders": [
                        placeholder.name for placeholder in slide_layout.placeholders
                    ],
                }
            )
        return {"available_layouts": layouts}

    def save_ppt(self, filename):
        """Save presentation with optional custom filename"""
        self.prs.save(filename)
        logger.info(f"Presentation saved as: {filename}")
