import json
import uuid
from typing import Any, Dict, List, Optional, Tuple

from pptx import Presentation

from libs.PPTMaker.platform.modules.bot.src.models.slide_models import ShapePosition
from libs.PPTMaker.platform.modules.bot.src.utils.charts_util import ChartHandler
from libs.PPTMaker.platform.modules.bot.src.utils.images_util import ImageHandler
from libs.PPTMaker.platform.modules.bot.src.utils.shapes_util import ShapeHandler
from libs.PPTMaker.platform.modules.bot.src.utils.slides_util import SlideLayoutManager
from libs.PPTMaker.platform.modules.bot.src.utils.styles.styling_util import (
    BasePresentationStyle,
)
from libs.PPTMaker.platform.modules.bot.src.utils.tables_util import TableHandler
from libs.utils.common.custom_logger import CustomLogger

log = CustomLogger("GroqLLMService", is_request=False)

logger, listener = log.get_logger()

listener.start()


class PPTXGenerator:
    """Main presentation generator class that orchestrates all components"""

    def __init__(self, style: BasePresentationStyle = None):
        """Initialize presentation generator with optional custom styling"""
        self.prs = Presentation()
        self.style = style or BasePresentationStyle()

        # Set slide dimensions
        self.prs.slide_width = self.style.get_dimension("slide_width")
        self.prs.slide_height = self.style.get_dimension("slide_height")

        # Initialize handlers
        self.layout_manager = SlideLayoutManager(self.prs, self.style)
        self.image_handler = ImageHandler(self.style)
        self.chart_handler = ChartHandler(self.style)
        self.table_handler = TableHandler(self.style)
        self.shape_handler = ShapeHandler(self.style)

    # Slide creation methods (delegate to appropriate handlers)
    def create_title_slide(self, title, subtitle=""):
        """Create a title slide"""
        return self.layout_manager.create_title_slide(title, subtitle)

    def create_content_slide(self, title, bullet_points):
        """Create a content slide with bullet points"""
        return self.layout_manager.create_content_slide(title, bullet_points)

    def add_image_slide(
        self, title, image_path, position=None, caption="", layout_type="standard"
    ):
        """Create an image slide"""
        return self.image_handler.add_image_slide(
            self.layout_manager, title, image_path, position, caption, layout_type
        )

    def add_chart_slide(self, title, chart_data, chart_type="column"):
        """Create a chart slide"""
        return self.chart_handler.add_chart_slide(
            self.layout_manager, title, chart_data, chart_type
        )

    def create_table_slide(self, title, table_data):
        """Create a table slide"""
        return self.table_handler.create_table_slide(
            self.layout_manager, title, table_data
        )

    def add_text_box(
        self, slide, text, left, top, width, height, font_style="body_large"
    ):
        """Add a text box to a slide"""
        return self.shape_handler.add_text_box(
            slide, text, left, top, width, height, font_style
        )

    def add_shape(self, slide, shape_type, left, top, width, height, fill_color=None):
        """Add a shape to a slide"""
        return self.shape_handler.add_shape(
            slide, shape_type, left, top, width, height, fill_color
        )

    def save_presentation(self, filename):
        """Save the presentation to a file"""
        try:
            self.prs.save(filename)
            print(f"Presentation saved as: {filename}")
            return True
        except Exception as e:
            print(f"Error saving presentation: {e}")
            return False


def create_sample_presentation():
    """Create a sample presentation demonstrating all features"""

    # Initialize generator
    ppt_gen = PPTXGenerator()

    # 1. Create title slide
    ppt_gen.create_title_slide(
        title="Python-PPTX Demo Presentation",
        subtitle="Demonstrating Basic PowerPoint Operations",
    )

    # 2. Create content slide with bullet points
    bullet_points = [
        "Create professional presentations programmatically",
        "Add various types of content (text, images, charts)",
        "Customize formatting and styling",
        "Automate repetitive presentation tasks",
        "Export to standard PowerPoint formats",
    ]
    ppt_gen.create_content_slide("Key Features", bullet_points)

    # 3. Create image slide (note: you'll need to provide an actual image path)
    # ppt_gen.add_image_slide(
    #     title="Sample Image Slide",
    #     image_path="sample_image.jpg",  # Replace with actual image path
    #     caption="This is a sample image caption"
    # )

    # 4. Create chart slide
    chart_data = {"Q1": 20, "Q2": 35, "Q3": 42, "Q4": 28}
    ppt_gen.add_chart_slide("Quarterly Results", chart_data, "column")

    # 5. Create table slide
    table_data = [
        ["Product", "Q1 Sales", "Q2 Sales", "Q3 Sales"],
        ["Product A", "$10,000", "$12,000", "$15,000"],
        ["Product B", "$8,000", "$9,500", "$11,000"],
        ["Product C", "$15,000", "$18,000", "$20,000"],
    ]
    ppt_gen.create_table_slide("Sales Data", table_data)

    # 6. Create a custom slide with shapes and text boxes
    blank_slide_layout = ppt_gen.prs.slide_layouts[6]
    custom_slide = ppt_gen.prs.slides.add_slide(blank_slide_layout)

    # Add custom title
    ppt_gen.add_text_box(custom_slide, "Custom Slide with Shapes", 1, 0.5, 11.33, 1, 28)

    # Add shapes
    ppt_gen.add_shape(custom_slide, "rectangle", 2, 2, 3, 1.5, (68, 114, 196))
    ppt_gen.add_shape(custom_slide, "circle", 8, 2, 2, 2, (255, 192, 0))

    # Add text over shapes
    ppt_gen.add_text_box(custom_slide, "Rectangle", 2.5, 2.5, 2, 0.5, 16)
    ppt_gen.add_text_box(custom_slide, "Circle", 8.5, 2.8, 1, 0.5, 16)

    # Save the presentation
    ppt_gen.save_presentation("sample_presentation.pptx")

    return ppt_gen


def create_quick_presentation(title, slides_data, filename):
    """Quickly create a presentation from structured data"""
    ppt_gen = PPTXGenerator()

    # Create title slide
    ppt_gen.create_title_slide(title)

    # Create content slides
    for slide_info in slides_data:
        if slide_info["type"] == "content":
            ppt_gen.create_content_slide(slide_info["title"], slide_info["content"])
        elif slide_info["type"] == "chart":
            ppt_gen.add_chart_slide(
                slide_info["title"],
                slide_info["data"],
                slide_info.get("chart_type", "column"),
            )
        elif slide_info["type"] == "table":
            ppt_gen.create_table_slide(slide_info["title"], slide_info["data"])

    # Save presentation
    ppt_gen.save_presentation(filename)
    return ppt_gen


# Example of quick presentation creation
if __name__ == "__main__":
    # Create sample presentation
    create_sample_presentation()

    # Example of quick presentation creation
    quick_slides = [
        {
            "type": "content",
            "title": "Project Overview",
            "content": [
                "Goal: Automate presentation creation",
                "Timeline: 4 weeks",
                "Team: 3 developers",
            ],
        },
        {
            "type": "chart",
            "title": "Progress Tracking",
            "data": {"Week 1": 25, "Week 2": 50, "Week 3": 75, "Week 4": 100},
            "chart_type": "line",
        },
    ]

    create_quick_presentation(
        "Project Status Update", quick_slides, "quick_presentation.pptx"
    )
