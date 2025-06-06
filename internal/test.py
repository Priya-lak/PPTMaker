import json
from pprint import pformat, pprint
from typing import Dict

from pptx import Presentation
from pptx.slide import Slide

from libs.PPTMaker.platform.modules.bot.src.constants import LAYOUT_GENERATION_PROMPT
from libs.PPTMaker.platform.modules.bot.src.models.slide_layout_models_2 import (
    LayoutDefinition,
    PresentationResponse,
)
from libs.PPTMaker.platform.modules.bot.src.services.llm_service import LLMService
from libs.PPTMaker.platform.modules.bot.src.services.ppt_generator_service import (
    PPTGenerator,
)
from libs.utils.common.custom_logger import CustomLogger

log = CustomLogger("TestLogger", is_request=False)

logger, listener = log.get_logger()
listener.start()


themes = [
    "static/templates/madison.pptx",
    "static/templates/nature.pptx",
    "static/templates/blue-spheres.pptx",
    "static/templates/ion-boardroom.pptx",
]


class PPT:
    def __init__(self, theme):
        self.prs = Presentation(theme)

    def get_theme_layouts(self):
        layouts = []
        for slide_layout in self.prs.slide_layouts:
            layout_blueprint = LayoutDefinition(
                layout=slide_layout.name,
                placeholders=[
                    {"idx": idx, "name": placeholder.name}
                    for idx, placeholder in enumerate(slide_layout.placeholders)
                ],
            )
            layouts.append(layout_blueprint.model_dump())

        logger.info(f"All layouts present in current theme {pformat(layouts)}")
        return layouts

    def add_slide(self, layout_identifier):
        """Add slide by layout name or index"""
        if isinstance(layout_identifier, str):
            # Find layout by name
            for i, layout in enumerate(self.prs.slide_layouts):
                if layout.name == layout_identifier:
                    # logger.info(f"adding slide {layout.name}")
                    slide = self.prs.slides.add_slide(self.prs.slide_layouts[i])
                    return slide
            raise ValueError(f"Layout '{layout_identifier}' not found")
        else:
            # Use as index
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[layout_identifier])
            return slide

    def add_sample_slides(self):
        for slide_layout in self.prs.slide_layouts:
            layout_index = self.prs.slide_layouts.index(slide_layout)
            slide = self.add_slide(layout_index)
            self.add_sample_text(slide)

    def populate_presentation(self, presentation_content: PresentationResponse):
        """Populate presentation with structured content"""
        logger.info("populating ppt")
        for slide_content in presentation_content.slides:
            try:
                # Add slide with the specified layout
                slide = self.add_slide(slide_content.layout)

                # Populate placeholders with content
                self.populate_slide_content(
                    slide, slide_content.placeholders, slide_content.layout
                )

            except ValueError as e:
                logger.info(
                    f"Error adding slide with layout '{slide_content.layout}': {e}"
                )
                continue
            except Exception as e:
                logger.info(f"Unexpected error processing slide: {e}")
                continue

    def populate_slide_content(
        self, slide: Slide, placeholder_content: Dict[str, str], sl_layout: str
    ):
        """Fill slide placeholders with provided content"""

        # Method 1: Direct iteration over slide placeholders
        original_slide_placeholders = [
            {"idx": idx, "name": pl.name} for idx, pl in enumerate(slide.placeholders)
        ]
        logger.debug(
            f"layout: {slide.slide_layout.name}, original slide placeholders {pformat(original_slide_placeholders)}"
        )
        logger.debug(f"layout: {sl_layout}, generated {pformat(placeholder_content)}")
        for idx, placeholder in enumerate(slide.placeholders):
            try:
                placeholder_name = placeholder.name
                content = placeholder_content.get(str(idx))
                if content is None:
                    continue

                # Check if placeholder can accept text
                if placeholder.has_text_frame:
                    if placeholder.text_frame is not None:
                        # Clear existing content
                        placeholder.text_frame.clear()

                        # Split content by newlines
                        lines = content.split("\\n")

                        # Add first line to the default paragraph
                        if lines:
                            placeholder.text_frame.text = lines[0]

                        # Add remaining lines as new paragraphs
                        for line in lines[1:]:
                            p = placeholder.text_frame.add_paragraph()
                            p.text = line

                        # logger.info(
                        #     f"Populated placeholder '{placeholder_name}' with content: {content}"
                        # )
                    else:
                        logger.warning(
                            f"Placeholder '{placeholder_name}' has text frame but frame is None"
                        )
                else:
                    logger.info(
                        f"Placeholder '{placeholder_name}' cannot accept text (no text frame)"
                    )

            except Exception as e:
                logger.error(
                    f"Error processing placeholder '{getattr(placeholder, 'name', 'unknown')}': {e}"
                )
                continue

    def get_available_layouts_json(self):
        """Get available layouts in JSON format for AI prompt"""
        layouts = []
        for slide_layout in self.prs.slide_layouts:
            layouts.append(
                {
                    "name": slide_layout.name,
                    "placeholders": [
                        placeholder.name for placeholder in slide_layout.placeholders
                    ],
                }
            )
        return {"available_layouts": layouts}

    def save_ppt(self, filename="output/demo/testing-layouts.pptx"):
        """Save presentation with optional custom filename"""
        self.prs.save(filename)
        logger.info(f"Presentation saved as: {filename}")


llm_service = LLMService()

service = PPT(theme=themes[-1])
layouts = service.get_theme_layouts()
messages = [
    {"role": "system", "content": LAYOUT_GENERATION_PROMPT},
    {"role": "system", "content": f"Main layout to generate content from: {layouts}"},
    {"role": "user", "content": "Topic: Machine learning algorithms"},
]
response = llm_service.call_llm(messages, {"type": "json_object"})


response = PresentationResponse(**json.loads(response))

service.populate_presentation(response)
# service.add_sample_slides()
service.save_ppt()

# content = {
#     "title": "The Magic of Photosynthesis: Unveiling the Power of Plant Life",
#     "presentation_content": [
#         {
#             "layout": "TITLE_AND_CONTENT",
#             "title": "The Magic of Photosynthesis: Unveiling the Power of Plant Life",
#             "points": [
#                 "Photosynthesis is like a magic power that plants use to make their own food from sunlight! It's a vital process that helps plants grow, and in turn, it provides oxygen for us to breathe.",
#                 "Without photosynthesis, our planet wouldn't be able to support life as we know it.",
#             ],
#         },
#         {
#             "layout": "TITLE_AND_CONTENT",
#             "title": "How Photosynthesis Works",
#             "points": [
#                 "photosynthesis is a two-step process.",
#                 "First, plants absorb water and nutrients from the ground through their roots.",
#                 "Then, they use sunlight, water, and a gas called carbon dioxide to create energy-rich sugars.",
#                 "This energy is stored in the plant's leaves, stems, and roots, allowing it to grow and thrive.",
#             ],
#         },
#         {
#             "layout": "TITLE_AND_CONTENT",
#             "title": "The Importance of Photosynthesis",
#             "points": [
#                 "photosynthesis is crucial for our planet's ecosystems.",
#                 "It provides oxygen for us to breathe, and it serves as the primary source of food for many animals.",
#                 "Additionally, plants help maintain healthy air and water quality by absorbing pollutants and excess carbon dioxide.",
#                 "We rely on photosynthesis every day, and it's incredible to think about the power and importance it holds!",
#             ],
#         },
#         {
#             "layout": "TITLE_AND_CONTENT",
#             "title": "Real-World Examples of Photosynthesis",
#             "points": [
#                 "Did you know that coral reefs, the Amazon rainforest, and even your backyard garden are all powered by photosynthesis?",
#                 "It's a reminder that photosynthesis is all around us, supporting life in every corner of our planet.",
#             ],
#         },
#         {"layout": "TITLE_ONLY", "title": "Conclusion"},
#         {
#             "layout": "TITLE_AND_CONTENT",
#             "title": "Conclusion",
#             "points": [
#                 "Understanding photosynthesis is crucial for appreciating the magic of plant life.",
#                 "By recognizing the importance of this process, we can begin to respect and care for the incredible plants that surround us.",
#                 "Join us as we continue to explore the wonders of photosynthesis and discover the incredible ways it impacts our world.",
#             ],
#         },
#     ],
# }


# layout = PresentationModel(**content)

# service = PPTGenerator(style=StylesEnum.CLASSY, theme=ThemesEnum.BLUE_SPHERES)
# service.create_presentation_from_layout(layout=layout)
